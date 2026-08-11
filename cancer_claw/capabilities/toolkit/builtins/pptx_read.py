

from __future__ import annotations

import asyncio
import base64
from pathlib import Path
from typing import Any

import structlog

from cancer_claw.capabilities.toolkit.base import BaseTool, ToolResult
from cancer_claw.capabilities.toolkit.workspace import resolve_tool_path, get_tool_workspace

logger = structlog.get_logger()

_SUPPORTED_IMAGE_CONTENT_TYPES = {
    "image/png", "image/jpeg", "image/gif", "image/webp",
    "image/bmp", "image/tiff",
}

_IMAGE_EXT_MAP = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/gif": ".gif",
    "image/webp": ".webp",
    "image/bmp": ".bmp",
    "image/tiff": ".tiff",
}

_SINGLE_IMAGE_MAX_BYTES = 20 * 1024 * 1024
_TOTAL_IMAGES_MAX = 30

class PptxReadTool(BaseTool):


    @property
    def name(self) -> str:
        return "pptx_read"

    @property
    def description(self) -> str:
        return (
            "读取 .pptx (PowerPoint) 文件的全部内容——逐幻灯片提取文本与嵌入图片。"
            "提取的图片会保存到 workspace 并以多模态方式注入上下文，让模型能同时"
            "阅读文字和查看图表/示意图。适用于会诊幻灯、学术报告、组会汇报等场景。"
        )

    def get_schema(self) -> dict:
        return {
            "type": "function",
            "function": {
                "name": "pptx_read",
                "description": self.description,
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": (
                                "PPT 文件路径（相对 workspace 或绝对路径）。"
                                "支持 .pptx 格式。"
                            ),
                        },
                        "include_images": {
                            "type": "boolean",
                            "description": (
                                "是否提取并注入图片到模型上下文（默认 true）。"
                                "设为 false 则只提取文本，节省 token。"
                            ),
                            "default": True,
                        },
                        "max_images": {
                            "type": "integer",
                            "description": (
                                f"最多提取的图片数量（默认 {_TOTAL_IMAGES_MAX}）。"
                                "超出后按幻灯片顺序截断。"
                            ),
                            "default": _TOTAL_IMAGES_MAX,
                        },
                        "slides": {
                            "type": "string",
                            "description": (
                                "要读取的幻灯片范围（可选）。"
                                "格式：'1-5' 或 '1,3,5,8-10'。"
                                "不传则读取全部幻灯片。"
                            ),
                        },
                    },
                    "required": ["path"],
                },
            },
        }

    async def execute(self, **kwargs: Any) -> ToolResult:
        raw_path = kwargs.get("path", "").strip()
        include_images = kwargs.get("include_images", True)
        max_images = min(kwargs.get("max_images", _TOTAL_IMAGES_MAX), _TOTAL_IMAGES_MAX)
        slides_spec = (kwargs.get("slides") or "").strip()

        if not raw_path:
            return ToolResult(success=False, error="path 参数不能为空")

        abs_path, err = resolve_tool_path(raw_path)
        if err:
            return ToolResult(success=False, error=f"路径解析失败: {err}")
        if not abs_path.exists():
            return ToolResult(success=False, error=f"文件不存在: {raw_path}")
        if abs_path.suffix.lower() not in (".pptx",):
            return ToolResult(
                success=False,
                error=f"不支持的文件格式 '{abs_path.suffix}'，仅支持 .pptx",
            )

        slide_indices = _parse_slide_spec(slides_spec) if slides_spec else None

        try:
            result = await asyncio.to_thread(
                _extract_pptx, abs_path, include_images, max_images, slide_indices
            )
        except ImportError:
            return ToolResult(
                success=False,
                error=(
                    "缺少 python-pptx 依赖。请在项目虚拟环境中执行：\n"
                    "  pip install python-pptx\n"
                    "然后重试。"
                ),
            )
        except Exception as e:
            logger.warning("pptx_read_failed", path=str(abs_path), error=str(e))
            return ToolResult(success=False, error=f"PPT 解析失败: {e}")

        text_output = result["text_output"]
        raw_images = result["images"]
        stats = result["stats"]


        summary_lines = [
            f"PPT 文件: {abs_path.name}",
            f"   总幻灯片: {stats['total_slides']} 页",
            f"   已读取: {stats['read_slides']} 页",
            f"   提取图片: {stats['extracted_images']} 张",
        ]
        if stats.get("skipped_images"):
            summary_lines.append(
                f"   跳过图片: {stats['skipped_images']} 张（超限/不支持格式）"
            )
        summary = "\n".join(summary_lines)
        output = f"{summary}\n\n{'═' * 60}\n\n{text_output}"


        data: dict[str, Any] = {"stats": stats}


        if raw_images and include_images:
            ws = get_tool_workspace()
            save_dir: Path | None = None
            if ws:
                save_dir = ws.default_relative_root / "pptx_images" / abs_path.stem
                save_dir.mkdir(parents=True, exist_ok=True)

            images_for_injection: list[dict[str, Any]] = []
            saved_paths: list[str] = []

            for img_info in raw_images:
                img_bytes: bytes = img_info["bytes"]
                content_type: str = img_info["content_type"]
                slide_num: int = img_info["slide_num"]
                img_idx: int = img_info["img_idx"]
                ext = _IMAGE_EXT_MAP.get(content_type, ".png")


                if save_dir:
                    filename = f"slide{slide_num:02d}_img{img_idx:02d}{ext}"
                    save_path = save_dir / filename
                    save_path.write_bytes(img_bytes)
                    saved_paths.append(f"pptx_images/{abs_path.stem}/{filename}")


                images_for_injection.append({
                    "bytes": img_bytes,
                    "mime": content_type,
                    "label": f"幻灯片{slide_num}-图{img_idx}",
                })

            data["images"] = images_for_injection
            if saved_paths:
                data["saved_image_paths"] = saved_paths
                output += (
                    f"\n\n图片已保存到 workspace/pptx_images/{abs_path.stem}/"
                    f"\n共 {len(saved_paths)} 张图片将注入视觉上下文供模型查看。"
                )

        return ToolResult(
            success=True,
            output=output,
            data=data,
        )

def _parse_slide_spec(spec: str) -> set[int]:

    indices: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            try:
                start, end = int(start_s.strip()), int(end_s.strip())
                indices.update(range(start, end + 1))
            except ValueError:
                continue
        else:
            try:
                indices.add(int(part))
            except ValueError:
                continue
    return indices

def _extract_pptx(
    path: Path,
    include_images: bool,
    max_images: int,
    slide_indices: set[int] | None,
) -> dict[str, Any]:

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    total_slides = len(prs.slides)
    text_parts: list[str] = []
    images: list[dict[str, Any]] = []
    extracted_count = 0
    skipped_count = 0
    read_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        if slide_indices and slide_num not in slide_indices:
            continue
        read_count += 1

        slide_texts: list[str] = []
        slide_texts.append(f"── 幻灯片 {slide_num} {'─' * 40}")


        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[备注] {notes_text}")

        for shape in slide.shapes:

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)


            if shape.has_table:
                table = shape.table
                table_lines = []
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    table_lines.append(" | ".join(cells))
                    if row_idx == 0:
                        table_lines.append(
                            "-" * (sum(len(c) for c in cells) + 3 * len(cells))
                        )
                if table_lines:
                    slide_texts.append("[表格]\n" + "\n".join(table_lines))


            if include_images and extracted_count < max_images:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ct = img.content_type
                    if ct in _SUPPORTED_IMAGE_CONTENT_TYPES:
                        blob = img.blob
                        if len(blob) <= _SINGLE_IMAGE_MAX_BYTES:
                            extracted_count += 1
                            images.append({
                                "bytes": blob,
                                "content_type": ct,
                                "slide_num": slide_num,
                                "img_idx": extracted_count,
                            })
                            slide_texts.append(
                                f"[图片 {extracted_count}: "
                                f"{img.filename or '嵌入图片'}]"
                            )
                        else:
                            skipped_count += 1
                            slide_texts.append(
                                f"[图片跳过: 超过大小限制 "
                                f"({len(blob) // 1024}KB)]"
                            )
                    else:
                        skipped_count += 1
                        slide_texts.append(f"[图片跳过: 不支持的格式 {ct}]")

        text_parts.append("\n".join(slide_texts))

    text_output = "\n\n".join(text_parts)
    return {
        "text_output": text_output,
        "images": images,
        "stats": {
            "total_slides": total_slides,
            "read_slides": read_count,
            "extracted_images": extracted_count,
            "skipped_images": skipped_count,
        },
    }
